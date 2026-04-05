"""
ppt_builder.py
将结构化 slide 数据渲染为 .pptx 文件。

支持：
  - 使用已有模板（读取首张 Slide 的主题色/字体）
  - 从零创建（内置 Ocean Gradient 配色）
  - 文本幻灯片、文本+表格幻灯片
  - 自动翻页（由 md_parser 处理）
"""

import io
import copy
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn

# ─────────────────────────────────────────────
# 默认主题配色 (Ocean Gradient)
# ─────────────────────────────────────────────
THEME = {
    "bg_dark":      RGBColor(0x06, 0x5A, 0x82),   # 深蓝 - 封面/结尾
    "bg_light":     RGBColor(0xF0, 0xF7, 0xFF),   # 淡蓝 - 内容页
    "accent":       RGBColor(0x02, 0xC3, 0x9A),   # 薄荷绿 - 强调
    "title_dark":   RGBColor(0xFF, 0xFF, 0xFF),   # 白色 - 深色背景标题
    "title_light":  RGBColor(0x06, 0x5A, 0x82),   # 深蓝 - 浅色背景标题
    "body":         RGBColor(0x1A, 0x1A, 0x2E),   # 深色正文
    "muted":        RGBColor(0x55, 0x77, 0x88),   # 次要文字
    "table_header": RGBColor(0x06, 0x5A, 0x82),
    "table_row1":   RGBColor(0xE8, 0xF4, 0xF8),
    "table_row2":   RGBColor(0xF8, 0xFC, 0xFF),
    "bullet_marker":RGBColor(0x02, 0xC3, 0x9A),
    "font_title":   "Calibri",
    "font_body":    "Calibri Light",
}

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ─────────────────────────────────────────────
# 公共入口
# ─────────────────────────────────────────────

def build_pptx(slides: list[dict], template_path: str | None = None) -> bytes:
    """
    将 slide 列表渲染为 pptx bytes。
    slides: md_parser.parse_markdown() 的输出
    template_path: 可选，已有 .pptx 模板路径
    """
    if template_path and Path(template_path).exists():
        prs = _load_template(template_path)
    else:
        prs = _create_blank_prs()

    # 清除模板中已有的幻灯片（保留 slide layouts）
    _clear_slides(prs)

    for i, slide_data in enumerate(slides):
        stype = slide_data.get("type", "content")
        if stype == "title" or (i == 0 and not slide_data.get("content")):
            _add_title_slide(prs, slide_data)
        elif stype == "long_text":
            _add_content_slide(prs, slide_data, is_long_text=True)
        else:
            _add_content_slide(prs, slide_data)

    # 序列化为 bytes
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ─────────────────────────────────────────────
# 模板处理
# ─────────────────────────────────────────────

def _load_template(path: str) -> Presentation:
    """加载模板，提取主题色替换全局 THEME。"""
    prs = Presentation(path)
    _extract_theme_from_template(prs)
    return prs


def _extract_theme_from_template(prs: Presentation):
    """从模板主题中尝试读取主色，部分更新 THEME。"""
    global THEME
    try:
        theme_element = prs.slide_master.element.find(
            './/' + qn('a:clrScheme'))
        if theme_element is None:
            return
        # 尝试读取 dk1 (深色1) 和 accent1
        dk1 = theme_element.find(qn('a:dk1'))
        accent1 = theme_element.find(qn('a:accent1'))

        def _read_color(el):
            if el is None:
                return None
            srgb = el.find('.//' + qn('a:srgbClr'))
            sys_clr = el.find('.//' + qn('a:sysClr'))
            if srgb is not None:
                val = srgb.get('val', '')
                if len(val) == 6:
                    return RGBColor.from_string(val)
            if sys_clr is not None:
                last = sys_clr.get('lastClr', '')
                if len(last) == 6:
                    return RGBColor.from_string(last)
            return None

        c1 = _read_color(dk1)
        c2 = _read_color(accent1)
        if c1:
            THEME["body"] = c1
            THEME["title_light"] = c1
        if c2:
            THEME["accent"] = c2
            THEME["bg_dark"] = c2
            THEME["table_header"] = c2
    except Exception:
        pass  # 模板解析失败则沿用默认主题


def _create_blank_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def _clear_slides(prs: Presentation):
    """删除 prs 中现有的所有幻灯片（保留 layouts/master）。"""
    # 倒序删除以避免索引问题（虽然此处用 slide_id 删）
    for i in range(len(prs.slides) - 1, -1, -1):
        _delete_slide_by_id(prs, prs.slides[i].slide_id)


def _delete_slide_by_id(prs: Presentation, slide_id: int):
    """按 slide_id 删除幻灯片。"""
    from pptx.oxml.ns import nsmap
    slides = prs.slides
    for i, slide in enumerate(slides):
        if slide.slide_id == slide_id:
            # 修复：prs 没有 presentation 属性，应直接使用 prs.element
            xml_slides = prs.element.find(qn('p:sldIdLst'))
            sldId_elements = xml_slides.findall(qn('p:sldId'))
            # 找到对应关系 ID
            rId = None
            for el in sldId_elements:
                if int(el.get('id')) == slide_id:
                    rId = el.get(qn('r:id'))
                    xml_slides.remove(el)
                    break
            if rId:
                prs.part.drop_rel(rId)
            break


# ─────────────────────────────────────────────
# 封面幻灯片
# ─────────────────────────────────────────────

def _add_title_slide(prs: Presentation, data: dict):
    layout = _get_blank_layout(prs)
    slide = prs.slides.add_slide(layout)
    w, h = prs.slide_width, prs.slide_height

    # 深色背景
    _fill_background(slide, THEME["bg_dark"])

    # 左侧装饰条
    bar = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0), Inches(0),
        Inches(0.5), h
    )
    _fill_shape(bar, THEME["accent"])
    bar.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(
        Inches(1.0), Inches(2.2),
        Inches(10), Inches(1.8)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = data.get("title", "")
    run.font.size = Pt(44)
    run.font.bold = True
    run.font.color.rgb = THEME["title_dark"]
    run.font.name = THEME["font_title"]
    p.alignment = PP_ALIGN.LEFT

    # 副标题
    subtitle = data.get("subtitle")
    if subtitle:
        sub_box = slide.shapes.add_textbox(
            Inches(1.0), Inches(4.2),
            Inches(10), Inches(1.0)
        )
        tf2 = sub_box.text_frame
        p2 = tf2.paragraphs[0]
        run2 = p2.add_run()
        run2.text = subtitle
        run2.font.size = Pt(22)
        run2.font.color.rgb = RGBColor(0xC0, 0xE0, 0xF0)
        run2.font.name = THEME["font_body"]

    # 底部色带
    footer_bar = slide.shapes.add_shape(
        1,
        Inches(0), h - Inches(0.6),
        w, Inches(0.6)
    )
    _fill_shape(footer_bar, THEME["accent"])
    footer_bar.line.fill.background()


# ─────────────────────────────────────────────
# 内容幻灯片（文本 / 表格 / 混合）
# ─────────────────────────────────────────────

def _add_content_slide(prs: Presentation, data: dict, is_long_text: bool = False):
    layout = _get_blank_layout(prs)
    slide = prs.slides.add_slide(layout)
    w, h = prs.slide_width, prs.slide_height

    # 浅色背景
    _fill_background(slide, THEME["bg_light"])

    # 顶部色带 + 标题
    header_h = Inches(1.1)
    header_bar = slide.shapes.add_shape(
        1,
        Inches(0), Inches(0),
        w, header_h
    )
    _fill_shape(header_bar, THEME["bg_dark"])
    header_bar.line.fill.background()

    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.1),
        w - Inches(1.0), header_h - Inches(0.15)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = data.get("title", "")
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = THEME["title_dark"]
    run.font.name = THEME["font_title"]
    p.alignment = PP_ALIGN.LEFT

    # 左侧细装饰条
    side_bar = slide.shapes.add_shape(
        1,
        Inches(0), header_h,
        Inches(0.08), h - header_h
    )
    _fill_shape(side_bar, THEME["accent"])
    side_bar.line.fill.background()

    # 内容区域
    content_top = header_h + Inches(0.25)
    content_left = Inches(0.35)
    content_w = w - Inches(0.7)
    content_h = h - content_top - Inches(0.3)

    content_blocks = data.get("content", [])
    _render_content_blocks(
        slide, content_blocks,
        content_left, content_top, content_w, content_h,
        is_long_text=is_long_text
    )


def _render_content_blocks(
    slide, blocks: list[dict],
    left, top, width, height,
    is_long_text=False
):
    """将 content blocks 渲染到幻灯片上。"""
    if not blocks:
        return

    # 分析块类型
    has_table = any(b["kind"] == "table" for b in blocks)

    if has_table:
        _render_mixed_content(slide, blocks, left, top, width, height)
    else:
        _render_text_content(slide, blocks, left, top, width, height, is_long_text)


def _render_text_content(slide, blocks, left, top, width, height, is_long_text):
    """纯文本/bullets 内容渲染。"""
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.word_wrap = True

    first_para = True
    for block in blocks:
        kind = block["kind"]

        if kind == "bullets":
            items = block["items"]
            is_ordered = block.get("ordered", False)
            for idx, item in enumerate(items):
                if first_para:
                    p = tf.paragraphs[0]
                    first_para = False
                else:
                    p = tf.add_paragraph()
                prefix = f"{idx + 1}. " if is_ordered else "• "
                run = p.add_run()
                run.text = prefix + item
                run.font.size = Pt(16) if not is_long_text else Pt(14)
                run.font.color.rgb = THEME["body"]
                run.font.name = THEME["font_body"]
                p.space_before = Pt(4)
                p.space_after = Pt(2)
                # 首字符颜色变化
                _color_bullet_prefix(run, THEME["accent"], len(prefix))

        elif kind == "text":
            if first_para:
                p = tf.paragraphs[0]
                first_para = False
            else:
                p = tf.add_paragraph()
            run = p.add_run()
            is_bold = block.get("bold", False)
            run.text = block["text"]
            run.font.size = Pt(15) if not is_long_text else Pt(13)
            run.font.bold = is_bold
            run.font.color.rgb = THEME["body"] if not is_bold else THEME["bg_dark"]
            run.font.name = THEME["font_body"]
            p.space_before = Pt(6)
            p.space_after = Pt(3)

        elif kind == "code":
            if first_para:
                p = tf.paragraphs[0]
                first_para = False
            else:
                p = tf.add_paragraph()
            run = p.add_run()
            run.text = block["text"]
            run.font.size = Pt(11)
            run.font.name = "Consolas"
            run.font.color.rgb = RGBColor(0x1E, 0x88, 0xE5)
            p.space_before = Pt(4)


def _color_bullet_prefix(run, color: RGBColor, prefix_len: int):
    """尝试对 bullet prefix 单独着色（通过 XML 操作）。"""
    # python-pptx 不支持在同一段内多个 run 颜色，
    # 这里通过添加额外 run 实现
    # 已经是单个 run，保持统一颜色
    pass


def _render_mixed_content(slide, blocks, left, top, width, height):
    """混合内容（文本 + 表格）渲染。"""
    # 分区：文本块在左，表格在右（或文本在上，表格在下）
    text_blocks = [b for b in blocks if b["kind"] != "table"]
    table_blocks = [b for b in blocks if b["kind"] == "table"]

    if text_blocks and table_blocks:
        # 左右布局：文本占 40%，表格占 58%
        text_w = width * 0.40
        table_left = left + text_w + Inches(0.2)
        table_w = width - text_w - Inches(0.2)

        if text_blocks:
            _render_text_content(slide, text_blocks, left, top, text_w, height)
        for tb in table_blocks:
            _render_table(slide, tb, table_left, top, table_w, height)
    else:
        # 纯表格
        cur_top = top
        for tb in table_blocks:
            tbl_h = min(height, Inches(0.4) * (len(tb.get("rows", [])) + 2))
            _render_table(slide, tb, left, cur_top, width, tbl_h)
            cur_top += tbl_h + Inches(0.2)


def _render_table(slide, table_block: dict, left, top, width, height):
    """渲染 Markdown 表格为 PPTX 表格。"""
    headers = table_block.get("headers", [])
    rows = table_block.get("rows", [])
    if not headers:
        return

    cols = len(headers)
    total_rows = 1 + len(rows)

    # 行高
    row_h = min(Inches(0.45), int(height / total_rows))

    tbl = slide.shapes.add_table(
        total_rows, cols,
        left, top,
        int(width), row_h * total_rows
    ).table

    # 列宽均分
    col_w = int(width / cols)
    for i in range(cols):
        tbl.columns[i].width = col_w

    # 表头行
    for ci, hdr in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.text = hdr
        tf = cell.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tf.paragraphs[0].runs[0] if tf.paragraphs[0].runs else tf.paragraphs[0].add_run()
        run.text = hdr
        run.font.bold = True
        run.font.size = Pt(13)
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        run.font.name = THEME["font_body"]
        _fill_cell(cell, THEME["table_header"])

    # 数据行
    for ri, row in enumerate(rows):
        bg = THEME["table_row1"] if ri % 2 == 0 else THEME["table_row2"]
        for ci in range(cols):
            cell = tbl.cell(ri + 1, ci)
            val = row[ci] if ci < len(row) else ""
            cell.text = val
            tf = cell.text_frame
            tf.word_wrap = True
            run = tf.paragraphs[0].runs[0] if tf.paragraphs[0].runs else tf.paragraphs[0].add_run()
            run.text = val
            run.font.size = Pt(12)
            run.font.color.rgb = THEME["body"]
            run.font.name = THEME["font_body"]
            _fill_cell(cell, bg)

    # 设置行高
    for ri in range(total_rows):
        tbl.rows[ri].height = row_h


# ─────────────────────────────────────────────
# 辅助函数
# ─────────────────────────────────────────────

def _get_blank_layout(prs: Presentation):
    """获取空白布局。"""
    for layout in prs.slide_layouts:
        if 'blank' in layout.name.lower() or layout.name == 'Blank':
            return layout
    # 退而求其次：返回最后一个或第一个
    return prs.slide_layouts[-1] if prs.slide_layouts else prs.slide_layouts[0]


def _fill_background(slide, color: RGBColor):
    """设置幻灯片背景色。"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _fill_shape(shape, color: RGBColor):
    """填充形状颜色。"""
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = color


def _fill_cell(cell, color: RGBColor):
    """填充表格单元格背景色。"""
    from pptx.oxml.ns import qn
    from lxml import etree
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    solidFill = etree.SubElement(tcPr, qn('a:solidFill'))
    srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
    srgbClr.set('val', str(color))
